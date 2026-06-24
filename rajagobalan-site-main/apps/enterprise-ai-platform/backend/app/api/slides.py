"""PowerPoint export endpoints."""

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from uuid import UUID
from datetime import datetime
import tempfile
import os

from app.database import get_db
from app.models import MaturityAssessment, DiscoveryResult, ROISimulation

router = APIRouter(prefix="/api/slides", tags=["slides"])


def create_presentation(company_id: UUID, db: Session):
    """Create a PowerPoint presentation with analysis results."""

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx not installed. Install with: pip install python-pptx"
        )

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)

    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "AI Transformation Strategy"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 2: Maturity Assessment
    assessment = (
        db.query(MaturityAssessment)
        .filter(MaturityAssessment.company_id == company_id)
        .order_by(MaturityAssessment.created_at.desc())
        .first()
    )

    if assessment:
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide2.shapes.title
        title.text = "AI Maturity Assessment"

        content = slide2.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        p = text_frame.paragraphs[0]
        p.text = f"Overall Score: {assessment.overall_score:.2f}/5.0"
        p.font.size = Pt(24)
        p.font.bold = True

        p = text_frame.add_paragraph()
        p.text = f"Maturity Level: {assessment.maturity_level.title()}"
        p.font.size = Pt(18)
        p.level = 0

        p = text_frame.add_paragraph()
        p.text = ""
        p.level = 0

        scores = [
            ("Strategy", assessment.strategy_score),
            ("Infrastructure", assessment.infrastructure_score),
            ("Data", assessment.data_score),
            ("People", assessment.people_score),
            ("Governance", assessment.governance_score),
        ]

        for domain, score in scores:
            p = text_frame.add_paragraph()
            p.text = f"{domain}: {score:.1f}/5.0"
            p.font.size = Pt(16)
            p.level = 1

    # Slide 3: Discovery Results
    discovery = (
        db.query(DiscoveryResult)
        .filter(DiscoveryResult.company_id == company_id)
        .order_by(DiscoveryResult.created_at.desc())
        .first()
    )

    if discovery:
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide3.shapes.title
        title.text = "AI Use Case Opportunities"

        content = slide3.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        p = text_frame.paragraphs[0]
        p.text = f"Industry: {discovery.industry}"
        p.font.size = Pt(16)

        p = text_frame.add_paragraph()
        p.text = "Top Opportunities:"
        p.font.size = Pt(14)
        p.font.bold = True

        opportunities = discovery.opportunities[:3]
        for opp in opportunities:
            p = text_frame.add_paragraph()
            p.text = f"{opp.get('use_case', '')}"
            p.font.size = Pt(12)
            p.level = 1

    # Slide 4: ROI Summary
    roi = (
        db.query(ROISimulation)
        .filter(ROISimulation.company_id == company_id)
        .order_by(ROISimulation.created_at.desc())
        .first()
    )

    if roi:
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide4.shapes.title
        title.text = "Expected ROI Impact"

        content = slide4.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        p = text_frame.paragraphs[0]
        p.text = f"ROI: {roi.roi_percentage:.1f}%"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(46, 158, 89)

        p = text_frame.add_paragraph()
        p.text = f"Net Benefit: ${roi.net_benefit:,.0f}"
        p.font.size = Pt(18)
        p.level = 0

        p = text_frame.add_paragraph()
        p.text = f"Payback Period: {roi.payback_months:.1f} months"
        p.font.size = Pt(16)
        p.level = 0

        p = text_frame.add_paragraph()
        p.text = f"NPV: ${roi.npv:,.0f}"
        p.font.size = Pt(16)
        p.level = 0

    # Slide 5: Next Steps
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    title.text = "Next Steps"

    content = slide5.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()

    steps = [
        "Establish AI governance and executive sponsorship",
        "Build or acquire AI/ML talent",
        "Develop data infrastructure and pipelines",
        "Execute pilot project to validate assumptions",
        "Scale successful pilots across the organization",
        "Establish AI Center of Excellence"
    ]

    for step in steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(14)
        p.level = 0

    return prs


@router.post("/export")
def export_presentation(
    company_id: UUID,
    db: Session = Depends(get_db)
):
    """Export analysis as PowerPoint presentation."""

    try:
        # Create presentation
        prs = create_presentation(company_id, db)

        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".pptx"
        )
        prs.save(temp_file.name)
        temp_file.close()

        # Return file
        return FileResponse(
            path=temp_file.name,
            filename=f"ai_strategy_{company_id}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
